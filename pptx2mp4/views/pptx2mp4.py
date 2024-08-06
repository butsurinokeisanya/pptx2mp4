
from flask import request, redirect, url_for, render_template, flash, send_file, session
from pptx2mp4 import app
from werkzeug.utils import secure_filename
import os
import cv2
import random
import string
import pyttsx3
import shutil

from pdf2image import convert_from_path
import glob
from moviepy.editor import *
import pptx
import time
from werkzeug.utils import secure_filename
import subprocess
import re
from pydub import AudioSegment


UPLOAD_FOLDER1 = r'.\pptx2mp4\static\pdf'
UPLOAD_FOLDER2 = r'.\pptx2mp4\static\pptx'
DOWNLOAD_FOLDER1=r'.\pptx2mp4\static\png'
DOWNLOAD_FOLDER2=r'.\pptx2mp4\static\wav'
DOWNLOAD_FOLDER3=r'.\pptx2mp4\static\mp4'
ALLOWED_EXTENSIONS1 = {'pdf'}
ALLOWED_EXTENSIONS2 = {'pptx'}


def allowed_file(filename, ALLOWED_EXTENSIONS):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def randomname(n):
    return ''.join(random.choices(string.ascii_letters + string.digits, k=n))

def pdf_to_images(pdf_path, output_folder):
    images = convert_from_path(pdf_path)
    os.makedirs(output_folder, exist_ok=True)
    for i, image in enumerate(images):
        image_path = os.path.join(output_folder, f"{i+1}.png")
        image.save(image_path, "PNG")

def convert_text_to_speech(text, output_file):
    engine = pyttsx3.init()
    engine.save_to_file(text, output_file)
    engine.runAndWait()

def process_pptx(input_pptx, output_folder):
    presentation = pptx.Presentation(input_pptx)
    os.makedirs("wav", exist_ok=True)
    for i, notes_slide in enumerate(presentation.slides, start=1):
        notes_text = ""
        if notes_slide.has_notes_slide and notes_slide.notes_slide:
            if notes_slide.notes_slide.notes_text_frame:
                for paragraph in notes_slide.notes_slide.notes_text_frame.paragraphs:
                    notes_text += paragraph.text + "\n"
            if notes_text:
                wav_filename = os.path.join(output_folder, f"{i}.wav")
                convert_text_to_speech(notes_text, wav_filename)
                print(f"Created {wav_filename} for Slide {i}")
                
def make_folder(folder_path):
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)

def delete_all_files_in_folder(folder_path):
    try:
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except Exception as e:
                print(f"Failed to delete {file_path}. Reason: {e}")
        print("All files and folders have been deleted successfully.")
    except Exception as e:
        print(f"Failed to access folder {folder_path}. Reason: {e}")

def sort_numerically(file_list):
    return sorted(file_list, key=lambda x: int(re.search(r'(\d+)', os.path.basename(x)).group(0)))


@app.route('/pptx2mp4_add', methods=['POST'])
def pptx2mp4_add():
    make_folder(UPLOAD_FOLDER1)
    make_folder(UPLOAD_FOLDER2)
    make_folder(DOWNLOAD_FOLDER1)
    make_folder(DOWNLOAD_FOLDER2)
    make_folder(DOWNLOAD_FOLDER3)
    delete_all_files_in_folder(UPLOAD_FOLDER1)
    delete_all_files_in_folder(UPLOAD_FOLDER2)
    delete_all_files_in_folder(DOWNLOAD_FOLDER1)
    delete_all_files_in_folder(DOWNLOAD_FOLDER2)
    delete_all_files_in_folder(DOWNLOAD_FOLDER3)
    image1 = request.files['upload_files1']
    if image1:
        if allowed_file(image1.filename, ALLOWED_EXTENSIONS1):
            image_name1 = secure_filename(image1.filename)
            image1.save(os.path.join(UPLOAD_FOLDER1, image_name1))
        else:
            image_name1 = None
            flash('pdfファイルではなかったので失敗しました.', 'alert alert-warning')
    else:
        image_name1 = None

    image2 = request.files['upload_files2']
    if image2:
        if allowed_file(image2.filename, ALLOWED_EXTENSIONS2):
            image_name2 = secure_filename(image2.filename)
            image2.save(os.path.join(UPLOAD_FOLDER2, image_name2))
        else:
            image_name2 = None
            flash('pptxファイルではなかったので失敗しました.', 'alert alert-warning')
    else:
        image_name2 = None

    pdf_path = os.path.join(UPLOAD_FOLDER1, image_name1)
    output_folder = DOWNLOAD_FOLDER1
    pdf_to_images(pdf_path, output_folder)

    pptx_path = os.path.join(UPLOAD_FOLDER2, image_name2)
    output_folder = DOWNLOAD_FOLDER2
    process_pptx(pptx_path, output_folder)

    png_files = glob.glob(f'{DOWNLOAD_FOLDER1}/*.png')
    png_files = sort_numerically(png_files)

    wav_files = glob.glob(f'{DOWNLOAD_FOLDER2}/*.wav')
    wav_files = sort_numerically(wav_files)

    if len(wav_files) != len(png_files):
        print("wavファイルとpngファイルの個数が等しくないので終了します.")
        sys.exit()

    clips = []
    for png_file, wav_file in zip(png_files, wav_files):
        audio_clip = AudioFileClip(wav_file)
        duration = audio_clip.duration
        image_clip = ImageClip(png_file).set_duration(duration)
        image_clip = image_clip.resize(newsize=(1600, 900))
        video_clip = CompositeVideoClip([image_clip.set_audio(audio_clip)])
        clips.append(video_clip)

    concat_clip = concatenate_videoclips(clips, method="compose")
    concat_clip.write_videofile(os.path.join(DOWNLOAD_FOLDER3, image_name1.split(".")[0] + ".mp4"), fps=24, write_logfile=True)

    flash('変換されました', 'alert alert-info')
    session['image_name1'] = image_name1  # Save image_name1 in the session
    return render_template('core/pptx2mp4.html',image_name1=image_name1)

@app.route('/pptx2mp4_download', methods=['GET'])
def pptx2mp4_download():
    image_name1 = session.get('image_name1')
    print(f"image_name1 retrieved from session: {image_name1}")  # Debugging statement
    if image_name1:
        return send_file(f'static/mp4/{image_name1.split(".")[0]}.mp4', as_attachment=True)
    else:
        flash('ファイルが見つかりませんでした。', 'alert alert-warning')
        return redirect(url_for('pptx2mp4_add'))



