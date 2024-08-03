from flask import request, redirect, url_for, render_template, flash, send_file
from pptx2mp4 import app
from werkzeug.utils import secure_filename
import os
import cv2
import random
import string
import shutil
from pdf2image import convert_from_path

import os
import pptx
import pyttsx3
import os
import subprocess
import re
from pydub import AudioSegment
import pptx


UPLOAD_FOLDER = r'.\pptx2mp4\static\pptx'
DOWNLOAD_FOLDER=r'.\pptx2mp4\static\wav'
ALLOWED_EXTENSIONS = {'pptx'}


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def randomname(n):
    return ''.join(random.choices(string.ascii_letters + string.digits, k=n))


def playVoicePeak(script, outpath, narrator="Japanese Female 4", happy=80, sad=0, angry=0, fun=10):
    exepath = "C:/Program Files/VOICEPEAK/voicepeak.exe"
    args = [
        exepath,
        "-s", script,
        "-n", narrator,
        "-o", outpath,
        "-e", f"happy={happy},sad={sad},angry={angry},fun={fun}"
    ]
    process = subprocess.Popen(args)
    process.communicate()

def split_script(script, max_length=140):
    sentences = re.split(r'(?<=。|！|\!|\.|\?|\？)', script)
    parts = []
    current_part = ""

    for sentence in sentences:
        if len(current_part) + len(sentence) <= max_length:
            current_part += sentence
        else:
            if current_part:
                parts.append(current_part)
            current_part = sentence

    if current_part:
        parts.append(current_part)

    return parts

def concatenate_wav_files(wav_files, output_path):
    combined = AudioSegment.empty()
    for wav_file in wav_files:
        audio_segment = AudioSegment.from_wav(wav_file)
        combined += audio_segment
    combined.export(output_path, format="wav")

def process_pptx(input_pptx, output_folder):
    presentation = pptx.Presentation(input_pptx)
    os.makedirs(output_folder, exist_ok=True)
    
    for i, slide in enumerate(presentation.slides, start=1):
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            if slide.notes_slide.notes_text_frame:
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                    notes_text += paragraph.text + "\n"
        
        if notes_text:
            scripts = split_script(notes_text)
            wav_files = []
            
            for j, script in enumerate(scripts):
                wav_filename = os.path.join(output_folder, f"{i}_{j}.wav")
                playVoicePeak(script, wav_filename)
                wav_files.append(wav_filename)
                print(f"Created {wav_filename} for Slide {i}, Part {j + 1}")
            
            final_wav_filename = os.path.join(output_folder, f"{i}.wav")
            concatenate_wav_files(wav_files, final_wav_filename)
            print(f"Concatenated wav file created at {final_wav_filename}")
            
            # 一時ファイルを削除
            for wav_file in wav_files:
                os.remove(wav_file)
                print(f"Deleted temporary file {wav_file}")
                
def delete_all_files_in_folder(folder_path):
    try:
        # フォルダの中のすべてのファイルとサブフォルダを削除
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)  # ファイルまたはリンクを削除
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)  # ディレクトリを削除
            except Exception as e:
                print(f"Failed to delete {file_path}. Reason: {e}")
        print("All files and folders have been deleted successfully.")
    except Exception as e:
        print(f"Failed to access folder {folder_path}. Reason: {e}")

@app.route('/txt2wav_voicepeak_add', methods=['POST'])
def txt2wav_voicepeak_add():
    image = request.files['upload_files']
    delete_all_files_in_folder(UPLOAD_FOLDER)
    if image:
        if allowed_file(image.filename):
            image_name = randomname(15) + '_' + secure_filename(image.filename)
            image.save(os.path.join(UPLOAD_FOLDER, image_name))
        else:
            image_name = None
            flash('pptxファイルではなかったので失敗しました.',
                  'alert alert-warning')
    else:
        image_name = None
    
    # 「CoInitialize は呼び出されていません。」で使えない.
    pdf_path=os.path.join(UPLOAD_FOLDER, image_name)
    output_folder=DOWNLOAD_FOLDER
    process_pptx(pdf_path, output_folder)

    flash('変換されました', 'alert alert-info')
    return render_template('core/txt2wav_voicepeak.html',image_name=image_name)

@app.route('/txt2wav_voicepeak_download',methods=['GET'])
def txt2wav_voicepeak_download():
    return send_file(r'static/png/output.pdf',as_attachment=True)

